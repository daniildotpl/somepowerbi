import pandas as pd
from ishch.connector import my_connection
from pptx import Presentation
from pptx.util import Inches, Pt



LIST_OF_TAGS = [
    'reporting_year', 
    'prevent_year', 
    'reporting_region', 
    'region_2024', 
    'kids_2024', 
    'adult_2024', 
    'count_org_2024', 
    'count_org_rank_2024', 
    'doctors_2024', 
    'doctors_rank_2024', 
    'nurse_2024', 
    'nurse_rank_2024', 
    'region_2025', 
    'kids_2025', 
    'adult_2025', 
    'count_org_2025', 
    'count_org_rank_2025', 
    'doctors_2025', 
    'doctors_rank_2025', 
    'nurse_2025', 
    'nurse_rank_2025', 
    'kids', 
    'adult'
]


def get_payload(engine, reporting_year, reporting_region):

    query = f"""
        WITH preset AS (
            SELECT 
                year,
                region,
                (SELECT COUNT(*) FROM statinfo 
                    WHERE year IN ({str(reporting_year)}, {str(reporting_year-1)}) AND region={reporting_region} AND tip_mo='kid polyclinic') kids, 
                (SELECT COUNT(*) FROM statinfo 
                    WHERE year IN ({str(reporting_year)}, {str(reporting_year-1)}) AND region={reporting_region} AND tip_mo='polyclinic') adult,
                COUNT(*) count_org,
                DENSE_RANK() OVER (PARTITION BY year ORDER BY COUNT(*) DESC) as count_org_rank,
                SUM(doctors) doctors,   
                DENSE_RANK() OVER (PARTITION BY year ORDER BY SUM(doctors) DESC) as doctors_rank,
                SUM(nurse) nurse,   
                DENSE_RANK() OVER (PARTITION BY year ORDER BY SUM(nurse) DESC) as nurse_rank
            FROM statinfo
            GROUP BY year, region
        )
        SELECT 
            year,
            region,
            kids,
            adult,
            count_org,
            count_org_rank,
            doctors,
            doctors_rank,
            nurse,
            nurse_rank
        FROM preset
        WHERE year IN ({str(reporting_year)}, {str(reporting_year-1)}) AND region={reporting_region}
        GROUP BY year, region, doctors, nurse
        ;
    """
    payload = pd.read_sql(query, engine)
    print(payload)
    return payload


def get_context(payload):
    context = {}
    context['reporting_year'] = payload.loc[1]['year']
    context['prevent_year'] = int(payload.loc[1]['year']) - 1
    context['reporting_region'] = payload.loc[1]['region']
    # ---
    payload_as_dict = payload.set_index('year').to_dict('index')
    for year, dictt in payload_as_dict.items():
        for key, value in dictt.items():
            key = key + '_' + year
            # print(key)
            context[key] = value
    # ---
    context['kids'] = payload.loc[1]['kids']
    context['adult'] = payload.loc[1]['adult']
    # print(f'context: {context}')
    return context


def make_powerp(context):

    prs = Presentation('ishch/template.pptx')  

    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                # print()
                # print('has_text_frame')
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        # print()
                        # print(f'run: {run.text}')
                        for tag in LIST_OF_TAGS:
                            if tag in run.text:
                                # print(f'find: {tag}')
                                try:
                                    value = context[tag]
                                    # print(f'context value: {value}')        
                                    run.text = run.text.replace(tag, str(context[tag]))
                                    # print('success')
                                except Exception as error:
                                    print(f'Ошибка в Run: {error}')
            if shape.has_table:
                # print()
                # print('has_table')
                table = shape.table
                for row in table.rows:
                    for cell in row.cells:
                        # print()
                        # print(f'cell: {cell.text}')
                        for tag in LIST_OF_TAGS:
                            if tag in cell.text:
                                # print(f'find: {tag}')
                                try:
                                    value = context[tag]
                                    # print(f'context value: {value}')        
                                    cell.text = cell.text.replace(tag, str(context[tag]))
                                    # print('success')
                                except Exception as error:
                                    print(f'Ошибка в Cell: {error}')

    prs.save('generated_presentation.pptx')


if __name__ == "__main__":
    payload = get_payload(my_connection(), 2025, 67)
    context = get_context(payload)
    make_powerp(context)













